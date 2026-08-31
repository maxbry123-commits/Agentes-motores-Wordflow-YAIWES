import base64
import json
from io import BytesIO
from pathlib import Path

import pandas as pd
from PIL import Image


def image_to_data_uri(image_path: str, output_format: str = "JPEG") -> str:
    path = Path(image_path)
    if not path.exists():
        raise FileNotFoundError(f"Image not found: {image_path}")

    image = Image.open(path)

    if output_format.upper() == "JPEG":
        if image.mode in ("RGBA", "LA"):
            background = Image.new("RGB", image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[-1])
            image = background
        elif image.mode == "P":
            image = image.convert("RGB")

    buffer = BytesIO()
    image.save(buffer, format=output_format)
    base64_str = base64.b64encode(buffer.getvalue()).decode("utf-8")
    return f"data:image/{output_format.lower()};base64,{base64_str}"


def transcribe_audio(file_path: str, client) -> str:
    with open(file_path, "rb") as audio_file:
        transcript = client.audio.transcriptions.create(
            model="gpt-4o-mini-transcribe", file=audio_file
        )

    return transcript.text


def load_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read().strip()


def load_excel(path: str) -> str:
    data = pd.read_excel(path, sheet_name=None, engine="openpyxl")
    return json.dumps(
        {sheet: df.to_dict(orient="records") for sheet, df in data.items()}, default=str
    )


def load_csv(path: str) -> str:
    df = pd.read_csv(path)
    return json.dumps(df.to_dict(orient="records"))


def load_docx(filepath: str) -> str:
    from docx import Document

    doc = Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs if p.text).strip()


def load_pptx(filepath: str) -> str:
    from pptx import Presentation

    prs = Presentation(filepath)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                chunks.append(shape.text)
    return "\n".join(s for s in chunks if s).strip()


def load_json(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def pdb_to_string(path: str) -> str:
    with open(path, "r") as file:
        content = file.read()
    return content


def load_pdf(filepath: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(filepath)
    text = " ".join(page.extract_text() or "" for page in reader.pages)
    return text.strip()
